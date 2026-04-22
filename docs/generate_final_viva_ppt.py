from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


TITLE_COLOR = RGBColor(12, 44, 84)
ACCENT_COLOR = RGBColor(0, 113, 188)
TEXT_COLOR = RGBColor(34, 34, 34)


def style_title(shape):
    p = shape.text_frame.paragraphs[0]
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    style_title(slide.shapes.title)

    tf = slide.placeholders[1].text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = TEXT_COLOR


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    first = True
    for level, text in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24 if level == 0 else 20)
            run.font.color.rgb = TEXT_COLOR


def add_two_content_slide(prs, title, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    left = slide.shapes.placeholders[1].text_frame
    left.clear()
    p = left.paragraphs[0]
    p.text = left_title
    p.level = 0
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.color.rgb = ACCENT_COLOR
    for item in left_points:
        q = left.add_paragraph()
        q.text = item
        q.level = 1
        q.runs[0].font.size = Pt(18)
        q.runs[0].font.color.rgb = TEXT_COLOR

    right = slide.shapes.placeholders[2].text_frame
    right.clear()
    p = right.paragraphs[0]
    p.text = right_title
    p.level = 0
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.color.rgb = ACCENT_COLOR
    for item in right_points:
        q = right.add_paragraph()
        q.text = item
        q.level = 1
        q.runs[0].font.size = Pt(18)
        q.runs[0].font.color.rgb = TEXT_COLOR


def add_table_slide(prs, title, columns, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    n_rows = len(rows) + 1
    n_cols = len(columns)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.8)
    )
    table = table_shape.table

    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.color.rgb = TITLE_COLOR

    for r, row_data in enumerate(rows, start=1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = TEXT_COLOR


def build_presentation(output_path):
    prs = Presentation()

    add_title_slide(
        prs,
        "Insider Threat Detections sytem",
        "Final Viva Presentation\nAayush Salve | Roll No: 22CC1007\nDept. of CSE Cybersecurity, RAIT | April 2026",
    )

    add_bullets_slide(
        prs,
        "Problem Statement",
        [
            (0, "Insider threats are hard to detect because users have valid access."),
            (0, "Rule-based SOC systems miss subtle misuse and create alert fatigue."),
            (0, "Need: intelligent, risk-prioritized, and explainable detection pipeline."),
        ],
    )

    add_bullets_slide(
        prs,
        "Project Objectives",
        [
            (0, "Detect unknown + known insider threats using hybrid ML."),
            (0, "Generate risk score (0-10) for analyst prioritization."),
            (0, "Reduce false alerts while preserving high threat coverage."),
            (0, "Deliver deployment-ready outputs: reports, API, dashboard."),
        ],
    )

    add_bullets_slide(
        prs,
        "System Architecture",
        [
            (0, "Input: login, file access, download, location, and activity logs."),
            (0, "Preprocessing + feature engineering pipeline."),
            (0, "IsolationForest (70%) for anomaly discovery."),
            (0, "RandomForest (30%) for threat validation."),
            (0, "Hybrid score -> Risk scoring engine -> SOC dashboard & alerts."),
        ],
    )

    add_two_content_slide(
        prs,
        "Technology Stack",
        "Core Stack",
        [
            "Python 3.12, Pandas, NumPy",
            "Scikit-learn, Joblib",
            "Flask backend + REST endpoints",
            "CSV-driven pipeline",
        ],
        "Visualization & Operations",
        [
            "Bootstrap UI",
            "Chart.js analytics",
            "Threat reports (CSV/JSON/TXT)",
            "SOC-friendly workflow",
        ],
    )

    add_bullets_slide(
        prs,
        "Dataset and Features",
        [
            (0, "Original dataset: 100 users (10 threat cases)."),
            (0, "Balanced dataset: 126 users (36 threat cases via SMOTE)."),
            (0, "Key features: failed logins, downloads, after-hours access,"),
            (1, "sensitive file access, unique access locations, anomaly score."),
        ],
    )

    add_bullets_slide(
        prs,
        "Modeling and Risk Logic",
        [
            (0, "Step 1: Train IsolationForest for unsupervised anomalies."),
            (0, "Step 2: Train RandomForest for supervised confirmation."),
            (0, "Step 3: Weighted hybrid score = 0.7*IF + 0.3*RF."),
            (0, "Step 4: Convert to risk score (0-10) and risk levels."),
            (0, "Threshold tuning performed for best F1."),
        ],
    )

    add_table_slide(
        prs,
        "Validation Results (5-Fold Cross-Validation)",
        ["Metric", "Average", "Std Dev", "Minimum", "Maximum"],
        [
            ["Precision", "0.88", "0.122", "0.67", "1.00"],
            ["Recall", "0.782", "0.160", "0.571", "1.00"],
            ["F1-Score", "0.82", "0.118", "0.615", "0.933"],
            ["Accuracy", "0.92", "0.030", "0.88", "0.96"],
        ],
    )

    add_bullets_slide(
        prs,
        "Threshold Optimization",
        [
            (0, "Optimal threshold selected: 0.5389"),
            (0, "At this threshold: Precision = 1.00, Recall = 1.00, AUC = 1.00"),
            (0, "Used to calibrate decision boundary in controlled evaluation."),
            (0, "Production monitoring still recommended for drift and recalibration."),
        ],
    )

    add_bullets_slide(
        prs,
        "Operational Threat Snapshot",
        [
            (0, "Latest report summary (100 users):"),
            (1, "Critical: 13, High: 51, Medium: 33, Low: 3"),
            (1, "Average risk score: 6.51 / 10"),
            (1, "Maximum risk score: 10.0"),
            (0, "Critical examples: USER_0077, USER_0064, USER_0065, USER_0037"),
        ],
    )

    add_bullets_slide(
        prs,
        "Key Findings",
        [
            (0, "Sensitive file access after hours is strongest indicator."),
            (0, "Failed logins and credential misuse are secondary indicators."),
            (0, "Multi-location access with high downloads is tertiary signal."),
            (0, "Hybrid design improves detection confidence vs single model."),
        ],
    )

    add_bullets_slide(
        prs,
        "Demo Flow for Viva",
        [
            (0, "1. Show dataset and engineered behavior features."),
            (0, "2. Run detection pipeline and generate risk scores."),
            (0, "3. Open dashboard: high-risk users, trend, and categories."),
            (0, "4. Explain one critical user case and recommended SOC action."),
        ],
    )

    add_two_content_slide(
        prs,
        "Limitations and Future Scope",
        "Current Limitations",
        [
            "Controlled dataset prototype",
            "Limited live SIEM integration",
            "No full forensic attribution layer",
            "Periodic retraining required",
        ],
        "Future Enhancements",
        [
            "Real-time streaming detection",
            "XAI-based analyst explanations",
            "Adaptive drift-triggered retraining",
            "Enterprise-scale deployment",
        ],
    )

    add_bullets_slide(
        prs,
        "Conclusion",
        [
            (0, "A practical hybrid ML framework for insider threat detection was built."),
            (0, "Model achieved strong validated performance (F1: 82%, Acc: 92%)."),
            (0, "System provides actionable risk scoring for SOC triage."),
            (0, "Project is deployment-ready with monitoring and retraining plan."),
        ],
    )

    add_bullets_slide(
        prs,
        "Thank You",
        [
            (0, "Questions and Discussion"),
            (0, "Guide: Dr. Dhananjay Dhakane"),
            (0, "Student: Aayush Salve (22CC1007)"),
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    output = r"d:\Major Project\Insider Threat Detection\docs\Insider_Threat_Detection_Final_Viva_Presentation.pptx"
    build_presentation(output)
    print(f"Presentation generated: {output}")
